from PIL import Image
from wand.image import Image as WandImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR 
import re
from os import listdir
from os.path import isfile, join

class presentation():
    def __init__(self, *args, **kwargs):
        self.i = 1

    def imageResize(self, img_path):
        img = Image.open(img_path)
        img.thumbnail((600, 400))
        img.save(img_path)
        return

    def watermark_images(self, image_path, image_folder, logo_path):       
        with WandImage(filename=image_path) as baseImage:
            width, height = baseImage.size
            logo_resolution = int(min(width, height)/3)
            logo = Image.open(logo_path)
            logo.thumbnail((logo_resolution, logo_resolution))
            logo.save(image_folder+r"/Watermark/logo_watermark.png")
            logo.close()

            with WandImage(filename=image_folder+r"/Watermark/logo_watermark.png") as water:
                baseImage.watermark(water, 0.0, 0, 0)
                baseImage.save(filename=r"./Images/Watermark/watermark_{}.jpg".format(self.i))
                self.i+=1
        return r"./Images/Watermark/watermark_{}.jpg".format(self.i-1)

    def add_image(self, slide, img_path):
        img = Image.open(img_path)
        width, height = img.size
        left = Inches(1.5)
        top = Inches(2.5)
        HEIGHT = Inches(4)
        WIDTH = Inches(7)
        if width >= height:
            pic = slide.shapes.add_picture(img_path, left, top, width=WIDTH)
        else:    
            pic = slide.shapes.add_picture(img_path, left, top, height=HEIGHT)

        return

    def add_title(self, slide):
        title = slide.shapes.add_textbox( left=Inches(1), 
                                        top=Inches(0.3), 
                                        height=Inches(1), 
                                        width=Inches(10))
        text_frame = title.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(44)
        font.bold = True
        font.italic = None 
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        run.text = 'Title'
        return slide

    def add_subtitle(self, slide):
        subtitle = slide.shapes.add_textbox( left=Inches(2), 
                                            top=Inches(1.5), 
                                            height=Inches(0.5), 
                                            width=Inches(8))
        text_frame = subtitle.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(32)
        font.bold = True
        font.italic = True  
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
        run.text = 'Subtitle'
        return slide

    def create_ppt(self, image_folder, logo_path):
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        images = [join(image_folder, f) for f in listdir(image_folder) if isfile(join(image_folder, f))]
        
        for img_path in images:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide = self.add_title(slide)
            slide = self.add_subtitle(slide)
            img_path = self.watermark_images(img_path)
            slide = self.add_image(slide, img_path)

        prs.save('test.pptx')


if __name__=="__main__":
    logo_path = r'D:\Trial_Projects_Python\watermark_ppt_assignment\nike_black.png'
    image_folder = r'D:\Trial_Projects_Python\watermark_ppt_assignment\Images'
    obj =  presentation().create_ppt(image_folder, logo_path)
    
